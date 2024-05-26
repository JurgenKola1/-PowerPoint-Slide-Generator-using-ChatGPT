import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI
from io import BytesIO
import requests
import config

# Initialize OpenAI client
client = OpenAI(api_key=config.API_KEY)

def generate_dalle_prompt(text):
    prompt = f"Summarize the following text to a DALL-E image generation prompt: \n{text}"
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": prompt}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )
    return response.choices[0].message.content

def generate_image(prompt):
    response = client.images.generate(
        model="dall-e-3",
        prompt=f"{prompt} Style: digital art",
        n=1,
        size="1024x1024"
    )
    return response.data[0].url

def generate_ppt_content(text, content_type):
    prompt = f"Create a {content_type} for a PowerPoint slide from the following text: \n{text}"
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": prompt}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    return response.choices[0].message.content

def add_slide(prs, text):
    dalle_prompt = generate_dalle_prompt(text)
    image_url = generate_image(dalle_prompt)
    bullet_points = generate_ppt_content(text, "bullet point text")
    slide_title = generate_ppt_content(text, "title")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))
    
    txBox = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(4), Inches(1.5))
    tf = txBox.text_frame
    tf.text = bullet_points
    
    title_shape = slide.shapes.title
    title_shape.text = slide_title

def get_slides():
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    prs.slide_width = Pt(1920)
    prs.slide_height = Pt(1080)
    
    for paragraph in paragraphs:
        add_slide(prs, paragraph)

    prs.save("my_presentation.pptx")

app = tk.Tk()
app.title("Create PPT Slides")
app.geometry("800x600")

text_field = tk.Text(app, wrap="word", font=("Arial", 12))
text_field.pack(fill="both", expand=True)
text_field.focus_set()

create_button = tk.Button(app, text="Create Slides", command=get_slides)
create_button.pack()

app.mainloop()

